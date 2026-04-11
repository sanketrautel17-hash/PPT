"""
Renderer — Phase 5
Template-preserving PPTX renderer using python-pptx.
Rules (non-negotiable):
- Only selected slides are kept; all others deleted in reverse index order.
- Placeholder text is replaced preserving run-level formatting.
- Chart data is replaced via replace_data only (type/style never changed).
- Images and decorative elements are NEVER touched.
- Obvious placeholder boilerplate text is cleaned.
"""

import io
import logging
from copy import deepcopy

from lxml import etree
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

from app.schemas.slide_plan import SlidePlan, SlidePlanItem

logger = logging.getLogger(__name__)

TEXTBOX_PLACEHOLDER_OFFSET = 10000

# Text to treat as placeholder boilerplate (cleaned if LLM leaves it)
BOILERPLATE_PATTERNS = [
    "insert", "click to add", "click to edit", "type here",
    "enter text", "add title", "add subtitle", "add text",
]
FILLER_TEXT_PATTERNS = [
    "lorem ipsum",
    "sed ut perspiciatis",
    "topic 1",
    "optional eyebrow",
    "subheadline",
    "eiludusponderium",
    "bullet point text",
    "presentation subtitle",
    "presenter name",
    "presentation title",
    "delete before use",
    "source: lorem ipsum",
]
NON_REPLACEABLE_TEXT_PATTERNS = [
    "confidential and proprietary",
    "unauthorized use, distribution, or reproduction prohibited",
]


def _is_boilerplate(text: str) -> bool:
    t = text.strip().lower()
    return any(pattern in t for pattern in BOILERPLATE_PATTERNS)

def _is_template_filler_text(text: str) -> bool:
    t = (text or "").strip().lower()
    if not t:
        return False
    if t.isdigit():
        return False
    if any(p in t for p in NON_REPLACEABLE_TEXT_PATTERNS):
        return False
    return any(p in t for p in FILLER_TEXT_PATTERNS)



def _is_compact_text_shape(shape) -> bool:
    """Detect genuinely tiny callout text boxes (e.g. radial-circle labels)."""
    try:
        width = int(shape.width or 0)
        height = int(shape.height or 0)
    except Exception:
        return False

    if width <= 0 or height <= 0:
        return False

    ratio = max(width, height) / max(1, min(width, height))

    # Tiny chips/circle labels.
    if width <= 950_000 and height <= 360_000:
        return True

    # Small near-compact labels.
    if width <= 1_100_000 and height <= 500_000 and ratio <= 1.9:
        return True

    return False


def _compact_label(value: str | list[str], shape) -> str:
    """Create a concise readable label for compact callouts based on shape size."""
    if isinstance(value, list):
        first = next((str(v).strip() for v in value if str(v).strip()), "")
    else:
        first = str(value or "").strip()

    if not first:
        return ""

    words = [w for w in first.replace("\n", " ").split() if w]
    if not words:
        return ""

    try:
        width = int(shape.width or 0)
        height = int(shape.height or 0)
    except Exception:
        width = 0
        height = 0

    # Extra-tiny chips: keep to one very short word.
    if width <= 550_000 or height <= 250_000:
        return words[0][:7]

    # Tiny circles: max one short word or a very short 2-word phrase.
    if width <= 900_000 and height <= 360_000:
        one = words[0][:10]
        if len(words) >= 2 and len(words[0]) + len(words[1]) + 1 <= 11:
            return f"{words[0]} {words[1]}"
        return one

    # Slightly larger compact labels: up to 2 words, trimmed.
    short = " ".join(words[:2]).strip()
    if len(short) <= 14:
        return short
    return words[0][:10]


def _replace_text_preserving_format(text_frame, new_text: str | list[str]) -> None:
    """
    Replace text in a TextFrame while preserving ALL run-level formatting.
    Handles both plain string and list-of-bullets.

    Each new bullet paragraph is styled using the corresponding original template
    paragraph at that position (pPr + first run), so per-level font sizes, indents,
    and bullet markers are preserved exactly. If there are more new bullets than
    original template paragraphs, the last template paragraph's style is reused.
    """
    from pptx.oxml.ns import qn

    if isinstance(new_text, list):
        bullets = [str(b) for b in new_text if str(b).strip()]
    else:
        bullets = [str(new_text).strip()]

    if not bullets:
        return

    txBody = text_frame._txBody

    # Collect existing paragraph XML nodes
    existing_paras = txBody.findall(qn("a:p"))

    # Build a per-paragraph template: (pPr_copy, run_copy) for each original paragraph.
    # This preserves per-level font sizes, indent markers, bullet chars, etc.
    para_templates: list[tuple] = []
    for para in existing_paras:
        pPr = para.find(qn("a:pPr"))
        pPr_copy = deepcopy(pPr) if pPr is not None else None
        run_copy = None
        for r in para.findall(qn("a:r")):
            run_copy = deepcopy(r)
            break
        para_templates.append((pPr_copy, run_copy))

    # Fallback when template had no paragraphs at all
    if not para_templates:
        para_templates = [(None, None)]

    # Remove ALL existing paragraph nodes from txBody
    for p in existing_paras:
        txBody.remove(p)

    # Re-build paragraphs — one per bullet, styled from the matching template paragraph.
    for i, bullet_text in enumerate(bullets):
        # Cap at last available template paragraph so extra bullets inherit its style.
        tmpl_idx = min(i, len(para_templates) - 1)
        pPr_tmpl, run_tmpl = para_templates[tmpl_idx]

        new_p = etree.SubElement(txBody, qn("a:p"))

        if pPr_tmpl is not None:
            new_p.insert(0, deepcopy(pPr_tmpl))

        if run_tmpl is not None:
            # Clone the template run to preserve font/size/bold/color
            run = deepcopy(run_tmpl)
            for t in run.findall(qn("a:t")):
                run.remove(t)
            t_elem = etree.SubElement(run, qn("a:t"))
            t_elem.text = bullet_text
            new_p.append(run)
        else:
            # No template run found — create a bare run
            run = etree.SubElement(new_p, qn("a:r"))
            t_elem = etree.SubElement(run, qn("a:t"))
            t_elem.text = bullet_text


def _replace_chart_data(shape, chart_data_model) -> None:
    """Replace chart data values only — never touch chart type or styling."""
    try:
        chart = shape.chart
        plot = chart.plots[0]

        from pptx.chart.data import ChartData as PptxChartData
        cd = PptxChartData()
        cd.categories = chart_data_model.categories

        for series in chart_data_model.series:
            cd.add_series(series.name, series.values)

        chart.replace_data(cd)
        logger.debug(f"Replaced chart data: {len(chart_data_model.series)} series, {len(chart_data_model.categories)} categories")
    except Exception as e:
        logger.warning(f"Chart data replacement failed: {e}")


def _render_slide(prs_slide, item: SlidePlanItem) -> None:
    """
    Fill placeholders and update chart data on a single slide.
    Images are never touched.
    """
    ph_content = item.content.placeholders
    slide_idx = item.template_slide_index

    # Log all placeholder keys available in plan vs what's on the slide
    plan_keys = set(ph_content.keys())
    slide_ph_keys = set()
    for shape_idx, shape in enumerate(prs_slide.shapes):
        if shape.is_placeholder:
            slide_ph_keys.add(str(shape.placeholder_format.idx))
        elif hasattr(shape, "text_frame") and shape.text_frame is not None:
            slide_ph_keys.add(str(TEXTBOX_PLACEHOLDER_OFFSET + shape_idx))
    matched = plan_keys & slide_ph_keys
    unmatched_plan = plan_keys - slide_ph_keys
    unmatched_slide = slide_ph_keys - plan_keys

    logger.info(
        f"[render] slide_index={slide_idx} | "
        f"plan_keys={sorted(plan_keys)} | "
        f"slide_ph_keys={sorted(slide_ph_keys)} | "
        f"matched={sorted(matched)} | "
        f"unmatched_in_plan={sorted(unmatched_plan)} | "
        f"missing_from_plan={sorted(unmatched_slide)}"
    )

    for shape_idx, shape in enumerate(prs_slide.shapes):
        try:
            # ── Placeholder text replacement ─────────────────────────────
            if shape.is_placeholder:
                ph = shape.placeholder_format
                ph_key = str(ph.idx)
                if ph_key in ph_content:
                    new_text = ph_content[ph_key]
                    if _is_compact_text_shape(shape):
                        new_text = _compact_label(new_text, shape)
                    preview = str(new_text[0] if isinstance(new_text, list) else new_text)
                    if new_text and not _is_boilerplate(preview):
                        _replace_text_preserving_format(shape.text_frame, new_text)
                        logger.info(
                            f"[render] slide={slide_idx} ph_idx={ph_key} "
                            f"→ wrote: '{preview[:60]}'"
                        )
                    else:
                        logger.warning(
                            f"[render] slide={slide_idx} ph_idx={ph_key} "
                            f"SKIPPED (empty or boilerplate): '{preview[:60]}'"
                        )
                else:
                    # LLM didn't provide content — clear leftover template filler
                    existing = shape.text_frame.text if shape.text_frame else ""
                    if _is_template_filler_text(existing):
                        _replace_text_preserving_format(shape.text_frame, "")
                        logger.info(
                            f"[render] slide={slide_idx} ph_idx={ph_key} "
                            f"CLEARED filler: '{existing.strip()[:60]}'"
                        )

            # ── Editable text-box replacement (pseudo-placeholder) ──────
            elif hasattr(shape, "text_frame") and shape.text_frame is not None:
                ph_key = str(TEXTBOX_PLACEHOLDER_OFFSET + shape_idx)
                if ph_key in ph_content:
                    new_text = ph_content[ph_key]
                    if _is_compact_text_shape(shape):
                        new_text = _compact_label(new_text, shape)
                    preview = str(new_text[0] if isinstance(new_text, list) else new_text)
                    if new_text and not _is_boilerplate(preview):
                        _replace_text_preserving_format(shape.text_frame, new_text)
                        logger.info(
                            f"[render] slide={slide_idx} pseudo_idx={ph_key} "
                            f"→ wrote: '{preview[:60]}'"
                        )
                    else:
                        logger.warning(
                            f"[render] slide={slide_idx} pseudo_idx={ph_key} "
                            f"SKIPPED (empty or boilerplate): '{preview[:60]}'"
                        )
                else:
                    # LLM didn't provide content — clear leftover template filler
                    existing = shape.text_frame.text if shape.text_frame else ""
                    if _is_template_filler_text(existing):
                        _replace_text_preserving_format(shape.text_frame, "")
                        logger.info(
                            f"[render] slide={slide_idx} pseudo_idx={ph_key} "
                            f"CLEARED filler: '{existing.strip()[:60]}'"
                        )

            # ── Chart data replacement ───────────────────────────────────
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART and item.content.chart_data:
                _replace_chart_data(shape, item.content.chart_data)

            # ── Pictures: NEVER TOUCH ────────────────────────────────────
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pass  # deliberately skip — images are sacred

        except Exception as e:
            logger.debug(f"Shape render error (slide {item.template_slide_index}, shape {shape.name}): {e}")
            continue


async def render_pptx(
    plan: SlidePlan,
    template_id: str,
    template_bytes: bytes,
) -> tuple[bytes, int]:
    """
    Render a complete PPTX from a SlidePlan and template binary.

    Steps:
    1. Open the original template.
    2. Build an ordered list of slide indices to keep.
    3. Delete all other slides in reverse index order.
    4. Fill placeholders and chart data on kept slides.
    5. Return the rendered PPTX as bytes.

    Args:
        plan: Validated SlidePlan from the LLM planner.
        template_id: ID string (used for logging only).
        template_bytes: Raw bytes of the master .pptx template.

    Returns:
        (pptx_bytes, slide_count)
    """
    prs = Presentation(io.BytesIO(template_bytes))
    total_template_slides = len(prs.slides)

    # Map template_slide_index -> SlidePlanItem for quick lookup.
    # Keep first occurrence and preserve plan order for final deck sequence.
    plan_map: dict[int, SlidePlanItem] = {}
    for item in plan.slides:
        if item.template_slide_index not in plan_map:
            plan_map[item.template_slide_index] = item
    keep_indices: list[int] = list(plan_map.keys())  # unique, first-occurrence order

    # Validate all indices are in range
    valid_keep = [i for i in keep_indices if 0 <= i < total_template_slides]
    if len(valid_keep) != len(keep_indices):
        invalid = [i for i in keep_indices if i < 0 or i >= total_template_slides]
        logger.warning(f"Template {template_id}: {invalid} are out of range — skipping")
        keep_indices = valid_keep

    # Identify slides to DELETE
    all_indices = set(range(total_template_slides))
    delete_indices = sorted(all_indices - set(keep_indices), reverse=True)  # reverse order to preserve index integrity

    # Delete unwanted slides (reverse order)
    xml_slides = prs.slides._sldIdLst
    slide_id_map = {i: prs.slides._sldIdLst[i] for i in range(total_template_slides)}

    for del_idx in delete_indices:
        try:
            slide_elem = prs.slides._sldIdLst[del_idx]
            rId = slide_elem.get("r:id")
            # Remove slide from XML
            xml_slides.remove(slide_elem)
            # Drop the relationship
            if rId and rId in prs.slides.part.rels:
                del prs.slides.part.rels[rId]
        except Exception as e:
            logger.debug(f"Could not delete slide {del_idx}: {e}")

    # Reorder remaining slides to match plan order (not template order).
    # This keeps narrative flow intact (e.g., Thank You must be last).
    ordered_elems = [slide_id_map[idx] for idx in keep_indices if idx in slide_id_map]
    current_elems = list(prs.slides._sldIdLst)
    for elem in current_elems:
        prs.slides._sldIdLst.remove(elem)
    for elem in ordered_elems:
        prs.slides._sldIdLst.append(elem)

    # Render content in the same order as keep_indices/plan order
    for slide_pos, slide in enumerate(prs.slides):
        if slide_pos >= len(keep_indices):
            continue
        orig_idx = keep_indices[slide_pos]
        item = plan_map.get(orig_idx)
        if item:
            _render_slide(slide, item)

    # Serialize to bytes
    output = io.BytesIO()
    prs.save(output)
    pptx_bytes = output.getvalue()

    slide_count = len(prs.slides)
    logger.info(f"Rendered PPTX for template {template_id}: {slide_count} slides, {len(pptx_bytes)} bytes")

    return pptx_bytes, slide_count

