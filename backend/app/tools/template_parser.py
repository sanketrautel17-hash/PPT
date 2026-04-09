"""
Template Parser — Phase 1
Extracts structural metadata from a PPTX binary using python-pptx + lxml.
Returns a fully populated TemplateProfile (slides, theme, placeholders, charts, images, tables).
"""

import io
import logging
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from app.schemas.template_profile import (
    BrandRules,
    ChartMeta,
    FontRules,
    ImageMeta,
    Placeholder,
    PlaceholderPosition,
    SlideLayout,
    TemplateProfile,
    TextStyle,
    Theme,
    ThemeColors,
    ThemeFonts,
    TableMeta,
)

logger = logging.getLogger(__name__)

# EMU per point (1 inch = 914400 EMU; 1 pt = 12700 EMU)
EMU_PER_PT = 12700

# Approximate characters per EMU of width at 1pt font
# Used to estimate max_chars: chars ≈ width_emu / (font_size_pt * EMU_PER_PT * 0.6)
CHAR_WIDTH_FACTOR = 0.6

# Placeholder type names (pptx enum → readable string)
PP_PLACEHOLDER_TYPES = {
    1: "TITLE",
    2: "BODY",
    3: "CENTER_TITLE",
    4: "SUBTITLE",
    5: "DATE",
    6: "SLIDE_NUMBER",
    7: "FOOTER",
    10: "PICTURE",
    13: "TABLE",
    14: "CHART",
    15: "SMART_ART",
    16: "MEDIA",
    18: "OBJECT",
    19: "SLIDE_IMAGE",
}


TEXTBOX_PLACEHOLDER_OFFSET = 10000
NON_EDITABLE_TEXT_PATTERNS = [
    "confidential and proprietary",
    "unauthorized use, distribution, or reproduction prohibited",
]
EDITABLE_TEXT_HINTS = [
    "lorem ipsum",
    "sed ut perspiciatis",
    "topic 1",
    "optional eyebrow",
    "presentation subtitle",
    "presenter name",
    "subheadline",
    "delete before use",
]
IGNORED_PLACEHOLDER_TYPES = {"DATE", "SLIDE_NUMBER", "FOOTER"}
GUIDANCE_TEXT_PATTERNS = [
    "read this info slide. this is not a template",
    "to download fonts",
    "brand guidelines",
    "a note on color",
    "guides",
    "grid system",
    "powerpoint icon library",
    "using slides templates",
    "familiarize yourself with the template",
    "start building your presentation",
    "follow brand guidelines",
    "proofread and review",
    "accessible presentations in powerpoint",
    "use the accessibility checker",
    "set reading order of slide contents",
    "add alt text to images",
    "do not merge or split cells in tables",
    "best practices checklist",
]


def _emu_position(shape) -> PlaceholderPosition:
    return PlaceholderPosition(
        left_emu=int(shape.left or 0),
        top_emu=int(shape.top or 0),
        width_emu=int(shape.width or 0),
        height_emu=int(shape.height or 0),
    )


def _estimate_max_chars(width_emu: int, font_size_pt: float, height_emu: int = 0) -> int:
    """
    Estimate max characters that can fit in a shape using both width and height.

    This is intentionally conservative for compact/square callouts (e.g. circle labels)
    to avoid overflow and clipped text in rendered decks.
    """
    if width_emu <= 0:
        return 200
    if not font_size_pt or font_size_pt <= 0:
        return 200

    safe_font_pt = font_size_pt
    chars_per_line = width_emu / (safe_font_pt * EMU_PER_PT * CHAR_WIDTH_FACTOR)

    if height_emu and height_emu > 0:
        line_height_emu = safe_font_pt * EMU_PER_PT * 1.35
        max_lines = max(1, int(height_emu / line_height_emu))
    else:
        max_lines = 2

    max_chars = int(chars_per_line * max_lines * 0.85)

    # Very short placeholders (e.g. eyebrow labels) should stay concise.
    if height_emu and height_emu < 400_000:
        max_chars = min(max_chars, 14)

    # Square/compact placeholders (often circular callouts) need extra-tight limits.
    if height_emu and min(width_emu, height_emu) > 0:
        ratio = max(width_emu, height_emu) / min(width_emu, height_emu)
        if ratio < 1.25 and max(width_emu, height_emu) < 2_200_000:
            max_chars = min(max_chars, 18)

    return max(4, max_chars)


def _fallback_font_size_pt(placeholder_type: str | None = None) -> float:
    p = (placeholder_type or "").upper()
    if p in {"TITLE", "CENTER_TITLE"}:
        return 32.0
    if p in {"SUBTITLE"}:
        return 24.0
    if p in {"BODY", "OBJECT", "CONTENT", "CENTER_BODY"}:
        return 18.0
    return 14.0


def _get_font_size(tf, placeholder_type: str | None = None) -> float:
    """Extract dominant font size from a text frame, with type-aware fallback."""
    try:
        for para in tf.paragraphs:
            for run in para.runs:
                if run.font.size:
                    return run.font.size / EMU_PER_PT
        # Try paragraph-level
        for para in tf.paragraphs:
            if para.runs and para.runs[0].font.size:
                return para.runs[0].font.size / EMU_PER_PT
    except Exception:
        pass
    return _fallback_font_size_pt(placeholder_type)


def _get_text_style(tf, shape) -> TextStyle:
    """Extract dominant text style from a text frame."""
    font_family = ""
    font_size_pt = 12.0
    bold = False
    italic = False
    color = ""

    try:
        for para in tf.paragraphs:
            for run in para.runs:
                f = run.font
                if f.name:
                    font_family = f.name
                if f.size:
                    font_size_pt = f.size / EMU_PER_PT
                if f.bold is not None:
                    bold = f.bold
                if f.italic is not None:
                    italic = f.italic
                if f.color and f.color.type is not None:
                    try:
                        color = f"#{f.color.rgb}"
                    except Exception:
                        color = "scheme:unknown"
                break
            if font_family:
                break
    except Exception:
        pass

    return TextStyle(
        font_family=font_family,
        font_size_pt=font_size_pt,
        bold=bold,
        italic=italic,
        color=color,
    )


def _parse_placeholder(shape, shape_idx: int) -> Placeholder | None:
    """Parse a placeholder shape into a Placeholder model."""
    try:
        ph = shape.placeholder_format
        if ph is None:
            return None

        ph_type_int = ph.type.real if ph.type else 2
        ph_type_name = PP_PLACEHOLDER_TYPES.get(ph_type_int, "BODY")
        if ph_type_name in IGNORED_PLACEHOLDER_TYPES:
            return None
        idx = ph.idx

        tf = shape.text_frame
        font_size_pt = _get_font_size(tf, ph_type_name)
        text_style = _get_text_style(tf, shape)
        current_text = tf.text.strip() if tf.text else ""
        if text_style.font_size_pt <= 0 or (not current_text and text_style.font_size_pt == 12.0):
            text_style.font_size_pt = font_size_pt

        position = _emu_position(shape)
        max_chars = _estimate_max_chars(position.width_emu, font_size_pt, position.height_emu)

        # Common template artifact: small numeric slide-number boxes detected as BODY placeholders
        if current_text.isdigit() and position.width_emu < 1_500_000 and position.height_emu < 600_000:
            return None

        return Placeholder(
            idx=idx,
            name=shape.name,
            type=ph_type_name,
            position=position,
            text_style=text_style,
            max_chars_estimate=max_chars,
            current_text=current_text,
        )
    except Exception as e:
        logger.debug(f"Could not parse placeholder shape {shape_idx}: {e}")
        return None





def _is_editable_text_candidate(text: str) -> bool:
    t = (text or "").strip().lower()
    if not t:
        return False
    if any(p in t for p in NON_EDITABLE_TEXT_PATTERNS):
        return False
    if t.isdigit():
        return False
    if any(h in t for h in EDITABLE_TEXT_HINTS):
        return True
    # Heuristic: generic short labels that templates expect authors to replace
    return t in {"text", "headline", "title", "subtitle", "source"}


def _classify_slide_texts(texts: list[str]) -> str:
    """Deterministically classify obvious guidance/instruction slides."""
    joined = "\n".join(t.lower() for t in texts if t)
    if any(p in joined for p in GUIDANCE_TEXT_PATTERNS):
        return "guidance"
    return "content_layout"


def _parse_textbox_as_placeholder(shape, shape_idx: int) -> Placeholder | None:
    """Treat editable non-placeholder text boxes as pseudo-placeholders."""
    try:
        if not hasattr(shape, "text_frame") or shape.text_frame is None:
            return None
        tf = shape.text_frame
        current_text = tf.text.strip() if tf.text else ""
        if not _is_editable_text_candidate(current_text):
            return None

        font_size_pt = _get_font_size(tf, "BODY")
        text_style = _get_text_style(tf, shape)
        if text_style.font_size_pt <= 0 or (not current_text and text_style.font_size_pt == 12.0):
            text_style.font_size_pt = font_size_pt

        position = _emu_position(shape)
        max_chars = _estimate_max_chars(position.width_emu, font_size_pt, position.height_emu)

        return Placeholder(
            idx=TEXTBOX_PLACEHOLDER_OFFSET + shape_idx,
            name=f"{shape.name} [text_box]",
            type="BODY",
            position=position,
            text_style=text_style,
            max_chars_estimate=max_chars,
            current_text=current_text,
        )
    except Exception as e:
        logger.debug(f"Could not parse text-box placeholder shape {shape_idx}: {e}")
        return None

def _parse_chart(shape, shape_idx: int) -> ChartMeta | None:
    """Parse a chart shape into ChartMeta."""
    try:
        chart = shape.chart
        chart_type = str(chart.chart_type).replace("XL_CHART_TYPE.", "")

        categories: list[str] = []
        series_names: list[str] = []
        current_values: list[list[float | None]] = []
        color_scheme: list[str] = []

        try:
            plot = chart.plots[0]
            # Categories
            if hasattr(plot, "series") and plot.series:
                first_ser = plot.series[0]
                try:
                    cat_labels = first_ser.data_labels
                except Exception:
                    cat_labels = None

                # Try to get category names from chart data
                try:
                    for xi, xval in enumerate(chart.plots[0].series[0].values):
                        categories.append(f"Cat{xi+1}")
                except Exception:
                    pass

                try:
                    categories = [
                        str(c) for c in chart.plots[0].series[0].data_labels
                    ]
                except Exception:
                    pass

            for ser in plot.series:
                series_names.append(ser.name or f"Series{len(series_names)+1}")
                vals: list[float | None] = []
                try:
                    vals = [v for v in ser.values]
                except Exception:
                    pass
                current_values.append(vals)
                if not categories and vals:
                    categories = [f"Cat{i+1}" for i in range(len(vals))]

        except Exception as e:
            logger.debug(f"Chart series parse error: {e}")

        return ChartMeta(
            shape_index=shape_idx,
            chart_type=chart_type,
            series_count=len(series_names),
            category_count=len(categories),
            categories=categories,
            series_names=series_names,
            current_values=current_values,
            color_scheme=color_scheme,
        )
    except Exception as e:
        logger.debug(f"Could not parse chart shape {shape_idx}: {e}")
        return None


def _parse_image(shape, shape_idx: int) -> ImageMeta:
    """Parse a picture/image shape into ImageMeta."""
    position = _emu_position(shape)
    # Heuristic: if at top-left corner and small → likely logo
    description = "logo" if (position.top_emu < 500_000 and position.width_emu < 2_000_000) else "image"
    return ImageMeta(shape_index=shape_idx, position=position, description=description)


def _parse_table(shape, shape_idx: int) -> TableMeta | None:
    """Parse a table shape into TableMeta."""
    try:
        table = shape.table
        rows = len(table.rows)
        cols = len(table.columns)
        header_row: list[str] = []
        if rows > 0:
            header_row = [cell.text.strip() for cell in table.rows[0].cells]

        # Try to extract font info from first cell
        cell_style: dict = {}
        try:
            cell = table.rows[0].cells[0]
            tf = cell.text_frame
            style = _get_text_style(tf, shape)
            cell_style = {"font_family": style.font_family, "font_size_pt": style.font_size_pt}
        except Exception:
            pass

        return TableMeta(
            shape_index=shape_idx,
            rows=rows,
            cols=cols,
            header_row=header_row,
            cell_style=cell_style,
        )
    except Exception as e:
        logger.debug(f"Could not parse table shape {shape_idx}: {e}")
        return None


# ─── Theme Extraction via lxml ───────────────────────────────────────────────

_SCHEME_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_THEME_ELEMENT_MAP = {
    "dk1": "dk1", "dk2": "dk2",
    "lt1": "lt1", "lt2": "lt2",
    "accent1": "accent1", "accent2": "accent2",
    "accent3": "accent3", "accent4": "accent4",
    "accent5": "accent5", "accent6": "accent6",
    "hlink": "hyperlink",
}


def _extract_theme(prs: Presentation) -> Theme:
    """Extract theme colors and fonts from ppt/theme/theme1.xml."""
    colors_dict: dict[str, str] = {}
    major_font = "Calibri"
    minor_font = "Calibri"

    try:
        # Access the theme part XML directly
        theme_part = prs.slide_master.theme_color_map
    except Exception:
        theme_part = None

    try:
        # Walk the prs package to find theme XML
        for rel in prs.slide_master.part.rels.values():
            if "theme" in rel.reltype:
                theme_xml = rel._target._blob
                root = ET.fromstring(theme_xml)
                ns = {"a": _SCHEME_NS}

                # Color scheme
                for color_elem in root.iter(f"{{{_SCHEME_NS}}}clrScheme"):
                    for child in color_elem:
                        tag = child.tag.split("}")[-1]  # strip namespace
                        mapped_key = _THEME_ELEMENT_MAP.get(tag, tag)
                        for sub in child:
                            val_tag = sub.tag.split("}")[-1]
                            if val_tag == "srgbClr":
                                val = sub.get("val", "")
                                colors_dict[mapped_key] = f"#{val.upper()}"
                            elif val_tag == "sysClr":
                                last_clr = sub.get("lastClr", "")
                                if last_clr:
                                    colors_dict[mapped_key] = f"#{last_clr.upper()}"

                # Font scheme
                for font_elem in root.iter(f"{{{_SCHEME_NS}}}fontScheme"):
                    for major_elem in font_elem.iter(f"{{{_SCHEME_NS}}}majorFont"):
                        latin = major_elem.find(f"{{{_SCHEME_NS}}}latin")
                        if latin is not None:
                            major_font = latin.get("typeface", major_font)
                    for minor_elem in font_elem.iter(f"{{{_SCHEME_NS}}}minorFont"):
                        latin = minor_elem.find(f"{{{_SCHEME_NS}}}latin")
                        if latin is not None:
                            minor_font = latin.get("typeface", minor_font)
                break
    except Exception as e:
        logger.debug(f"Theme extraction error: {e}")

    theme_colors = ThemeColors(**{k: v for k, v in colors_dict.items() if k in ThemeColors.model_fields})
    return Theme(
        colors=theme_colors,
        fonts=ThemeFonts(major_font=major_font, minor_font=minor_font),
    )


# ─── Main Parser ─────────────────────────────────────────────────────────────

async def parse_template(pptx_bytes: bytes, template_name: str) -> TemplateProfile:
    """
    Parse a PPTX binary and return a complete TemplateProfile.

    Args:
        pptx_bytes: Raw bytes of the .pptx file.
        template_name: Display name of the template.

    Returns:
        TemplateProfile with all slides, theme, and metadata populated.
    """
    prs = Presentation(io.BytesIO(pptx_bytes))
    theme = _extract_theme(prs)

    slide_layouts: list[SlideLayout] = []

    for slide_idx, slide in enumerate(prs.slides):
        placeholders: list[Placeholder] = []
        charts: list[ChartMeta] = []
        images: list[ImageMeta] = []
        tables: list[TableMeta] = []
        slide_text_samples: list[str] = []

        for shape_idx, shape in enumerate(slide.shapes):
            try:
                if hasattr(shape, "text_frame") and shape.text_frame is not None:
                    raw_text = shape.text_frame.text.strip() if shape.text_frame.text else ""
                    if raw_text:
                        slide_text_samples.append(raw_text)

                # ── Placeholder ──────────────────────
                if shape.is_placeholder:
                    ph = _parse_placeholder(shape, shape_idx)
                    if ph:
                        placeholders.append(ph)

                # ── Editable text box (pseudo-placeholder) ───────────────
                elif hasattr(shape, "text_frame") and shape.text_frame is not None:
                    ph = _parse_textbox_as_placeholder(shape, shape_idx)
                    if ph:
                        placeholders.append(ph)

                # ── Chart ────────────────────────────
                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    chart = _parse_chart(shape, shape_idx)
                    if chart:
                        charts.append(chart)

                # ── Picture / Image ──────────────────
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    images.append(_parse_image(shape, shape_idx))

                # ── Table ────────────────────────────
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = _parse_table(shape, shape_idx)
                    if table:
                        tables.append(table)

                # ── Group (search inside) ────────────
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for sub_idx, sub_shape in enumerate(shape.shapes):
                        if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            images.append(_parse_image(sub_shape, shape_idx * 100 + sub_idx))

            except Exception as e:
                logger.debug(f"Slide {slide_idx}, shape {shape_idx} parse error: {e}")
                continue

        slide_layouts.append(
            SlideLayout(
                slide_index=slide_idx,
                classification=_classify_slide_texts(slide_text_samples),
                placeholders=placeholders,
                charts=charts,
                images=images,
                tables=tables,
            )
        )

    profile = TemplateProfile(
        name=template_name,
        theme=theme,
        slides=slide_layouts,
        total_slides=len(slide_layouts),
        usable_layouts=len(slide_layouts),  # updated after guidance extraction
    )

    logger.info(
        f"Parsed template '{template_name}': {profile.total_slides} slides, "
        f"theme fonts={theme.fonts.major_font}/{theme.fonts.minor_font}"
    )
    return profile

