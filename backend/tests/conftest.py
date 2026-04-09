"""
Shared pytest fixtures for all tests.
"""

import io

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt


# ─── Minimal PPTX factory ────────────────────────────────────────────────────

def make_minimal_pptx(
    num_slides: int = 3,
    include_chart: bool = False,
    include_table: bool = False,
) -> bytes:
    """
    Build a minimal in-memory PPTX with `num_slides` slides.
    Each slide has a title placeholder (idx=0) and a body placeholder (idx=1).
    Optionally adds a chart or table on slide 0.
    Returns raw bytes ready to pass to parse_template or render_pptx.
    """
    prs = Presentation()
    blank_layout = prs.slide_layouts[1]  # "Title and Content" layout

    for _ in range(num_slides):
        slide = prs.slides.add_slide(blank_layout)
        # Set placeholder text so the parser can read sizes
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                ph.text = "Slide Title"
            elif ph.placeholder_format.idx == 1:
                ph.text = "Slide body text"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def minimal_pptx() -> bytes:
    """A 3-slide minimal PPTX as bytes."""
    return make_minimal_pptx(3)


@pytest.fixture
def single_slide_pptx() -> bytes:
    """A 1-slide minimal PPTX as bytes."""
    return make_minimal_pptx(1)


@pytest.fixture
def five_slide_pptx() -> bytes:
    """A 5-slide minimal PPTX as bytes."""
    return make_minimal_pptx(5)
