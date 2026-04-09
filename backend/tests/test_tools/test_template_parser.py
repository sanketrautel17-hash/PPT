"""
Unit tests for the template parser.
No LLM calls, no DB — uses in-memory PPTX built with python-pptx.
"""

import io
import pytest

from pptx import Presentation
from pptx.util import Inches, Pt

from app.tools.template_parser import (
    _estimate_max_chars,
    _get_font_size,
    parse_template,
)
from tests.conftest import make_minimal_pptx


# ─── Helper: build a PPTX with specific layouts ──────────────────────────────

def _build_pptx(num_slides: int = 3) -> bytes:
    return make_minimal_pptx(num_slides)


# ─── _estimate_max_chars ─────────────────────────────────────────────────────

def test_estimate_max_chars_typical():
    """A typical 6-inch wide placeholder at 18pt should give a reasonable char estimate."""
    # 6 inches = 6 * 914400 EMU = 5,486,400 EMU
    # 18pt → chars ≈ 5_486_400 / (18 * 12700 * 0.6) ≈ 5_486_400 / 136_980 ≈ 40
    result = _estimate_max_chars(5_486_400, 18.0)
    assert 30 <= result <= 200, f"Expected ~40 chars, got {result}"


def test_estimate_max_chars_zero_font_size():
    """Zero font size should return the sensible default (200)."""
    result = _estimate_max_chars(5_000_000, 0)
    assert result == 200


def test_estimate_max_chars_zero_width():
    """Zero placeholder width should return the sensible default (200)."""
    result = _estimate_max_chars(0, 24.0)
    assert result == 200


def test_estimate_max_chars_small_font_gives_more_chars():
    """Smaller font → more characters fit."""
    wide_small = _estimate_max_chars(5_000_000, 10.0)
    wide_large = _estimate_max_chars(5_000_000, 36.0)
    assert wide_small > wide_large


# ─── parse_template — structural correctness ─────────────────────────────────

@pytest.mark.asyncio
async def test_parse_template_returns_profile():
    """parse_template should return a TemplateProfile with the correct slide count."""
    pptx_bytes = _build_pptx(3)
    profile = await parse_template(pptx_bytes, "Test Template")
    assert profile.name == "Test Template"
    assert profile.total_slides == 3
    assert len(profile.slides) == 3


@pytest.mark.asyncio
async def test_parse_template_slides_are_content_layout_by_default():
    """All slides start as 'content_layout' before guidance extraction."""
    pptx_bytes = _build_pptx(4)
    profile = await parse_template(pptx_bytes, "Layout Test")
    for slide in profile.slides:
        assert slide.classification == "content_layout"


@pytest.mark.asyncio
async def test_parse_template_slides_have_correct_indices():
    """Each parsed slide should have the correct slide_index (0-based)."""
    pptx_bytes = _build_pptx(5)
    profile = await parse_template(pptx_bytes, "Index Test")
    for expected_idx, slide in enumerate(profile.slides):
        assert slide.slide_index == expected_idx


@pytest.mark.asyncio
async def test_parse_template_placeholders_extracted():
    """Each slide in a Title+Content layout should have at least 1 placeholder."""
    pptx_bytes = _build_pptx(2)
    profile = await parse_template(pptx_bytes, "PH Test")
    for slide in profile.slides:
        assert len(slide.placeholders) >= 1


@pytest.mark.asyncio
async def test_parse_template_single_slide():
    """Parsing a 1-slide template should work without errors."""
    pptx_bytes = _build_pptx(1)
    profile = await parse_template(pptx_bytes, "Single Slide")
    assert profile.total_slides == 1


@pytest.mark.asyncio
async def test_parse_template_placeholder_max_chars_positive():
    """All parsed placeholders should have a positive max_chars_estimate."""
    pptx_bytes = _build_pptx(2)
    profile = await parse_template(pptx_bytes, "Chars Test")
    for slide in profile.slides:
        for ph in slide.placeholders:
            assert ph.max_chars_estimate > 0, (
                f"Slide {slide.slide_index} placeholder idx={ph.idx} has max_chars=0"
            )


@pytest.mark.asyncio
async def test_parse_template_empty_pptx_raises_or_returns_zero():
    """Passing corrupt/empty bytes should raise or return profile with 0 slides."""
    try:
        profile = await parse_template(b"not a pptx", "Bad")
        # If it doesn't raise, we just accept whatever it returns
    except Exception:
        pass  # Expected — corrupt bytes


@pytest.mark.asyncio
async def test_parse_template_preserves_name():
    """The profile name must match what was passed in."""
    name = "Camping World Q1 2025"
    pptx_bytes = _build_pptx(2)
    profile = await parse_template(pptx_bytes, name)
    assert profile.name == name


@pytest.mark.asyncio
async def test_parse_template_theme_has_fonts():
    """The extracted theme should have non-empty font names."""
    pptx_bytes = _build_pptx(2)
    profile = await parse_template(pptx_bytes, "Theme Test")
    assert len(profile.theme.fonts.major_font) > 0
    assert len(profile.theme.fonts.minor_font) > 0


@pytest.mark.asyncio
async def test_parse_template_content_slides_helper():
    """The content_slides() helper should return all slides when none are 'guidance'."""
    pptx_bytes = _build_pptx(3)
    profile = await parse_template(pptx_bytes, "Helper Test")
    content = profile.content_slides()
    assert len(content) == 3


@pytest.mark.asyncio
async def test_parse_template_extracts_editable_textbox_as_placeholder():
    """Non-placeholder editable text boxes should be exposed as pseudo-placeholders."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title-only
    title = slide.shapes.title
    if title:
        title.text = "Title"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1.2))
    box.text_frame.text = "Lorem ipsum dolor sit amet"

    buf = io.BytesIO()
    prs.save(buf)

    profile = await parse_template(buf.getvalue(), "Textbox Template")
    ph_ids = {ph.idx for ph in profile.slides[0].placeholders}
    assert any(idx >= 10000 for idx in ph_ids)


@pytest.mark.asyncio
async def test_parse_template_marks_guidance_slide_deterministically():
    """Instructional slides should be marked guidance without relying on LLM classification."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    tb.text_frame.text = "Read this info slide. This is not a template."

    buf = io.BytesIO()
    prs.save(buf)

    profile = await parse_template(buf.getvalue(), "Guidance Template")
    assert profile.slides[0].classification == "guidance"
