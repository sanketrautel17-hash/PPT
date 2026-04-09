"""
Unit tests for the PPTX renderer.
No LLM calls, no DB — uses in-memory PPTX built with python-pptx.
"""

import io
import pytest

from pptx import Presentation
from pptx.util import Inches

from app.schemas.slide_plan import ChartData, ChartSeries, SlideContent, SlidePlan, SlidePlanItem
from app.tools.renderer import _is_boilerplate, _replace_text_preserving_format, render_pptx
from tests.conftest import make_minimal_pptx


# ─── _is_boilerplate ─────────────────────────────────────────────────────────

@pytest.mark.parametrize("text,expected", [
    ("Click to add title", True),
    ("click to edit master title style", True),
    ("insert subtitle", True),
    ("type here", True),
    ("add subtitle", True),
    ("Q1 Business Review", False),
    ("Revenue Growth Summary", False),
    ("", False),
    ("   ", False),
])
def test_is_boilerplate(text: str, expected: bool):
    assert _is_boilerplate(text) is expected


# ─── render_pptx — basic smoke test ──────────────────────────────────────────

@pytest.mark.asyncio
async def test_render_pptx_returns_bytes():
    """render_pptx should return non-empty bytes and a positive slide count."""
    pptx_bytes = make_minimal_pptx(3)
    plan = SlidePlan(slides=[
        SlidePlanItem(
            slide_type="title",
            template_slide_index=0,
            purpose="Title slide",
            content=SlideContent(placeholders={"0": "Q1 Review", "1": "Finance Team"}),
        ),
    ])
    result_bytes, slide_count = await render_pptx(plan, "test_id", pptx_bytes)
    assert isinstance(result_bytes, bytes)
    assert len(result_bytes) > 0
    assert slide_count >= 1


@pytest.mark.asyncio
async def test_render_pptx_slide_count_matches_plan():
    """The rendered PPTX should have as many slides as requested in the plan."""
    pptx_bytes = make_minimal_pptx(5)
    plan = SlidePlan(slides=[
        SlidePlanItem(
            slide_type=f"slide_{i}",
            template_slide_index=i,
            purpose=f"Slide {i}",
            content=SlideContent(placeholders={"0": f"Title {i}"}),
        )
        for i in range(3)  # request 3 out of 5
    ])
    result_bytes, slide_count = await render_pptx(plan, "test_id", pptx_bytes)
    assert slide_count == 3


@pytest.mark.asyncio
async def test_render_pptx_output_is_valid_pptx():
    """Rendered bytes should be openable as a valid PPTX."""
    pptx_bytes = make_minimal_pptx(3)
    plan = SlidePlan(slides=[
        SlidePlanItem(
            slide_type="title",
            template_slide_index=0,
            purpose="Title",
            content=SlideContent(placeholders={"0": "Hello World"}),
        )
    ])
    result_bytes, _ = await render_pptx(plan, "test", pptx_bytes)
    # Should open without raising
    prs = Presentation(io.BytesIO(result_bytes))
    assert len(prs.slides) >= 1


@pytest.mark.asyncio
async def test_render_pptx_with_bullet_list():
    """render_pptx should handle bullet list content without raising."""
    pptx_bytes = make_minimal_pptx(3)
    plan = SlidePlan(slides=[
        SlidePlanItem(
            slide_type="bullets",
            template_slide_index=1,
            purpose="Key points",
            content=SlideContent(
                placeholders={
                    "0": "Key Points",
                    "1": ["Revenue up 20%", "Customer retention 95%", "New markets entered"],
                }
            ),
        )
    ])
    result_bytes, slide_count = await render_pptx(plan, "test_bullets", pptx_bytes)
    assert isinstance(result_bytes, bytes)
    assert slide_count == 1


@pytest.mark.asyncio
async def test_render_pptx_out_of_range_indices_skipped():
    """Slide indices beyond the template range should be silently skipped."""
    pptx_bytes = make_minimal_pptx(3)  # only indices 0-2 valid
    plan = SlidePlan(slides=[
        SlidePlanItem(
            slide_type="slide",
            template_slide_index=0,
            purpose="Valid slide",
            content=SlideContent(placeholders={"0": "Title"}),
        ),
        SlidePlanItem(
            slide_type="slide",
            template_slide_index=99,  # out of range
            purpose="Invalid slide",
            content=SlideContent(placeholders={"0": "Should be skipped"}),
        ),
    ])
    result_bytes, slide_count = await render_pptx(plan, "test_oob", pptx_bytes)
    assert isinstance(result_bytes, bytes)
    # Only 1 valid slide should be kept
    assert slide_count == 1


@pytest.mark.asyncio
async def test_render_pptx_all_slides_selected():
    """Selecting all slides from the template should keep them all."""
    pptx_bytes = make_minimal_pptx(3)
    plan = SlidePlan(slides=[
        SlidePlanItem(
            slide_type=f"slide_{i}",
            template_slide_index=i,
            purpose=f"Slide {i}",
            content=SlideContent(placeholders={"0": f"Title {i}"}),
        )
        for i in range(3)
    ])
    result_bytes, slide_count = await render_pptx(plan, "all_slides", pptx_bytes)
    assert slide_count == 3


@pytest.mark.asyncio
async def test_render_pptx_empty_plan_results_in_zero_slides():
    """An empty slide plan should result in a PPTX with 0 slides kept."""
    pptx_bytes = make_minimal_pptx(3)
    plan = SlidePlan(slides=[])
    result_bytes, slide_count = await render_pptx(plan, "empty", pptx_bytes)
    assert isinstance(result_bytes, bytes)
    assert slide_count == 0


@pytest.mark.asyncio
async def test_render_pptx_boilerplate_text_is_not_written():
    """
    If the plan content starts with boilerplate text, the renderer skips writing it.
    The slide title should remain unchanged (empty or template default).
    """
    pptx_bytes = make_minimal_pptx(3)
    plan = SlidePlan(slides=[
        SlidePlanItem(
            slide_type="title",
            template_slide_index=0,
            purpose="Boilerplate test",
            content=SlideContent(placeholders={"0": "Click to add title"}),
        )
    ])
    result_bytes, _ = await render_pptx(plan, "bplate", pptx_bytes)
    prs = Presentation(io.BytesIO(result_bytes))
    slide = prs.slides[0]
    # The title placeholder should NOT have been updated with the boilerplate text
    title_ph = None
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 0:
            title_ph = shape
            break
    if title_ph:
        assert "click to add" not in title_ph.text_frame.text.lower()


@pytest.mark.asyncio
async def test_render_pptx_writes_pseudo_textbox_placeholder():
    """Renderer should fill editable non-placeholder text boxes via synthetic keys."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title-only
    if slide.shapes.title:
        slide.shapes.title.text = "Template Title"
    body_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(7), Inches(1.2))
    body_box.text_frame.text = "Lorem ipsum dolor sit amet"

    template_shape_idx = None
    for i, shp in enumerate(slide.shapes):
        if hasattr(shp, "text_frame") and shp.text_frame is not None:
            if (shp.text_frame.text or "").strip() == "Lorem ipsum dolor sit amet":
                template_shape_idx = i
                break
    assert template_shape_idx is not None
    pseudo_key = str(10000 + template_shape_idx)

    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    plan = SlidePlan(slides=[
        SlidePlanItem(
            slide_type="body",
            template_slide_index=0,
            purpose="Body content",
            content=SlideContent(placeholders={pseudo_key: "Generated body content"}),
        )
    ])

    result_bytes, _ = await render_pptx(plan, "pseudo", pptx_bytes)
    out = Presentation(io.BytesIO(result_bytes))
    out_texts = []
    for shp in out.slides[0].shapes:
        if hasattr(shp, "text_frame") and shp.text_frame is not None:
            out_texts.append((shp.text_frame.text or "").strip())

    assert "Generated body content" in out_texts


@pytest.mark.asyncio
async def test_render_pptx_does_not_autofill_unmapped_textboxes():
    """Renderer should only write mapped placeholders to avoid accidental overlap text."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    if slide.shapes.title:
        slide.shapes.title.text = "Template Title"
    body_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(7), Inches(1.2))
    body_box.text_frame.text = "Lorem ipsum dolor sit amet"

    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    plan = SlidePlan(slides=[
        SlidePlanItem(
            slide_type="body",
            template_slide_index=0,
            purpose="Revenue improved across key segments",
            content=SlideContent(placeholders={"0": "Revenue Overview"}),
        )
    ])

    result_bytes, _ = await render_pptx(plan, "no_autofill", pptx_bytes)
    out = Presentation(io.BytesIO(result_bytes))
    texts = []
    for shp in out.slides[0].shapes:
        if hasattr(shp, "text_frame") and shp.text_frame is not None:
            t = (shp.text_frame.text or "").strip()
            if t:
                texts.append(t)

    assert "Revenue improved across key segments" not in texts
    assert any("lorem ipsum" in t.lower() for t in texts)


@pytest.mark.asyncio
async def test_render_pptx_respects_plan_slide_order():
    """Slides should be ordered according to plan, not original template indices."""
    pptx_bytes = make_minimal_pptx(4)
    plan = SlidePlan(slides=[
        SlidePlanItem(
            slide_type="middle_first",
            template_slide_index=2,
            purpose="First output slide",
            content=SlideContent(placeholders={"0": "Planned First"}),
        ),
        SlidePlanItem(
            slide_type="opening_second",
            template_slide_index=0,
            purpose="Second output slide",
            content=SlideContent(placeholders={"0": "Planned Second"}),
        ),
    ])

    result_bytes, slide_count = await render_pptx(plan, "order", pptx_bytes)
    assert slide_count == 2

    out = Presentation(io.BytesIO(result_bytes))
    first_title = out.slides[0].shapes.title.text.strip() if out.slides[0].shapes.title else ""
    second_title = out.slides[1].shapes.title.text.strip() if out.slides[1].shapes.title else ""

    assert first_title == "Planned First"
    assert second_title == "Planned Second"


@pytest.mark.asyncio
async def test_render_pptx_compact_textbox_is_compressed():
    """Long text in tiny callouts should be auto-compressed to short labels."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tiny = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(0.35), Inches(0.2))
    tiny.text_frame.text = "Topic 1"

    tiny_idx = None
    for i, shp in enumerate(slide.shapes):
        if hasattr(shp, "text_frame") and shp.text_frame is not None:
            if (shp.text_frame.text or "").strip() == "Topic 1":
                tiny_idx = i
                break
    assert tiny_idx is not None
    pseudo_key = str(10000 + tiny_idx)

    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    plan = SlidePlan(slides=[
        SlidePlanItem(
            slide_type="compact",
            template_slide_index=0,
            purpose="Compact callout",
            content=SlideContent(placeholders={pseudo_key: "Higher Service Revenue"}),
        )
    ])

    result_bytes, _ = await render_pptx(plan, "compact", pptx_bytes)
    out = Presentation(io.BytesIO(result_bytes))

    rendered = None
    for shp in out.slides[0].shapes:
        if hasattr(shp, "text_frame") and shp.text_frame is not None:
            t = (shp.text_frame.text or "").strip()
            if t and t != "Topic 1":
                rendered = t
                break

    assert rendered is not None
    assert len(rendered) <= 11


@pytest.mark.asyncio
async def test_render_pptx_wide_short_textbox_not_compacted():
    """Wide short textboxes (e.g. source lines) should not be compacted to acronyms."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    source_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(2.8), Inches(0.2))
    source_box.text_frame.text = "Source: Internal Analysis"

    src_idx = None
    for i, shp in enumerate(slide.shapes):
        if hasattr(shp, "text_frame") and shp.text_frame is not None:
            if (shp.text_frame.text or "").strip() == "Source: Internal Analysis":
                src_idx = i
                break
    assert src_idx is not None

    pseudo_key = str(10000 + src_idx)
    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    plan = SlidePlan(slides=[
        SlidePlanItem(
            slide_type="source",
            template_slide_index=0,
            purpose="Source text",
            content=SlideContent(placeholders={pseudo_key: "Source: Internal Sales Data"}),
        )
    ])

    result_bytes, _ = await render_pptx(plan, "source", pptx_bytes)
    out = Presentation(io.BytesIO(result_bytes))

    texts = []
    for shp in out.slides[0].shapes:
        if hasattr(shp, "text_frame") and shp.text_frame is not None:
            t = (shp.text_frame.text or "").strip()
            if t:
                texts.append(t)

    assert "Source: Internal Sales Data" in texts

